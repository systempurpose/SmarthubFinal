from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple

from docx import Document
from pptx import Presentation


DOCX_PATH = Path("docx/newly-updated21.docx")
OUT_PPTX_PATH = Path("pptx/proposal-defense.pptx")


@dataclass(frozen=True)
class SlideSpec:
    title: str
    sources: Sequence[re.Pattern]
    max_bullets: int = 8


def normalize_space(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def is_all_caps_heading(text: str) -> bool:
    t = normalize_space(text)
    if not t or len(t) > 70:
        return False
    has_letters = bool(re.search(r"[A-Za-z]", t))
    return has_letters and t.upper() == t


def is_marker(text: str) -> bool:
    t = normalize_space(text)
    if not t:
        return False

    if re.match(r"^(CHAPTER|Chapter)\s+\d+\b", t):
        return True

    # numbered headings, including ones that start with a dot (seen in some docs)
    if re.match(r"^\.?\d+(?:\.\d+)*\b", t):
        return True

    if is_all_caps_heading(t):
        return True

    # Common section titles (not necessarily numbered)
    common = [
        "Project Context",
        "Purpose and Description",
        "Scope and Delimitation",
        "Definition of Terms",
        "Objectives of The Project",
        "Review of related Literature",
        "Review of related Literature and Studies",
        "Conceptual Framework",
    ]
    if any(t.lower() == c.lower() for c in common):
        return True

    return False


def extract_nonempty_paragraphs(doc: Document) -> List[str]:
    out: List[str] = []
    for p in doc.paragraphs:
        t = normalize_space(p.text)
        if t:
            out.append(t)
    return out


def compute_marker_indices(lines: Sequence[str]) -> List[int]:
    return [i for i, t in enumerate(lines) if is_marker(t)]


def find_heading_index(lines: Sequence[str], patterns: Sequence[re.Pattern]) -> Optional[int]:
    # Prefer exact match first
    for i, t in enumerate(lines):
        for pat in patterns:
            if pat.fullmatch(t):
                return i

    # Then allow search
    for i, t in enumerate(lines):
        for pat in patterns:
            if pat.search(t):
                return i

    return None


def slice_to_next_marker(lines: Sequence[str], start_index: int, marker_indices: Sequence[int]) -> List[str]:
    # Content begins after the heading line
    start = start_index + 1
    end = len(lines)
    for mi in marker_indices:
        if mi > start_index:
            end = mi
            break
    return list(lines[start:end])


def clean_bullet(text: str) -> str:
    t = normalize_space(text)
    t = t.replace("�", "-")

    # Strip common bullet/numbering prefixes
    t = re.sub(r"^\s*[•\-*]+\s+", "", t)
    t = re.sub(r"^\s*\(?\d+(?:\.\d+)*\)?\s*[-.)]?\s*", "", t)
    return t.strip()


def to_bullets(paragraphs: Iterable[str], limit: int) -> List[str]:
    bullets: List[str] = []
    for raw in paragraphs:
        t = clean_bullet(raw)
        if not t:
            continue
        # Skip chapter labels
        if re.match(r"^(CHAPTER|Chapter)\s+\d+\b", t, re.I):
            continue
        if t.upper() == t and len(t) <= 40 and re.search(r"[A-Z]", t):
            # likely a header line repeated in content
            continue
        bullets.append(t)
        if len(bullets) >= limit:
            break
    return bullets


def split_bullets(bullets: Sequence[str], first_n: int) -> Tuple[List[str], List[str]]:
    return list(bullets[:first_n]), list(bullets[first_n:])


def extract_project_title(lines: Sequence[str]) -> str:
    # Prefer the canonical project title string, if present anywhere.
    canonical = re.compile(
        r"\bSmart\s*Hub\s*:\s*Diagnostic\s+Application\s+for\s+Android\s+Phone\b",
        re.I,
    )
    for t in lines:
        if canonical.search(t):
            return "Smart Hub: Diagnostic Application for Android Phone"

    # Fallback: first reasonable-looking mention.
    for t in lines:
        if "smart hub" in t.lower() and ":" in t:
            m = re.search(r"(Smart\s*Hub\s*:\s*[^.\"]{10,80})", t, re.I)
            if m:
                return normalize_space(m.group(1))
    return "Proposal Defense"


def add_title_slide(prs: Presentation, title: str, members: Sequence[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    subtitle = slide.placeholders[1]
    if members:
        subtitle.text = "Members:\n" + "\n".join(members)
    else:
        subtitle.text = "Members:\n(please add member names)"


def add_bullets_slide(prs: Presentation, title: str, bullets: Sequence[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    if not bullets:
        p = body.paragraphs[0]
        p.text = "(content not found in DOCX — please verify section text)"
        return

    first = True
    for b in bullets:
        if first:
            p = body.paragraphs[0]
            p.text = b
            first = False
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0


def main() -> int:
    if not DOCX_PATH.exists():
        raise SystemExit(f"DOCX not found: {DOCX_PATH}")

    doc = Document(str(DOCX_PATH))
    lines = extract_nonempty_paragraphs(doc)
    markers = compute_marker_indices(lines)

    # Extract intro chunk so we can split it into two slides.
    intro_idx = find_heading_index(lines, [re.compile(r"INTRODUCTION", re.I)])
    intro_paras = slice_to_next_marker(lines, intro_idx, markers) if intro_idx is not None else []

    # Objectives
    obj_idx = find_heading_index(lines, [re.compile(r"Objectives of The Project", re.I)])
    obj_paras = slice_to_next_marker(lines, obj_idx, markers) if obj_idx is not None else []

    # Scope
    scope_idx = find_heading_index(lines, [re.compile(r"Scope and Delimitation", re.I)])
    scope_paras = slice_to_next_marker(lines, scope_idx, markers) if scope_idx is not None else []

    # Review of Literature & Studies (broad section)
    rol_idx = find_heading_index(lines, [re.compile(r"REVIEW OF RELATED LITERATURE AND STUDIES", re.I)])
    rol_paras = slice_to_next_marker(lines, rol_idx, markers) if rol_idx is not None else []

    # Summary of RRL (explicit)
    rrl_sum_idx = find_heading_index(lines, [re.compile(r"2\.2\.6\s+Summary of Review Related Literature and Studies", re.I)])
    rrl_sum_paras = slice_to_next_marker(lines, rrl_sum_idx, markers) if rrl_sum_idx is not None else []

    # Conceptual framework
    cf_idx = find_heading_index(lines, [re.compile(r"Conceptual Framework", re.I)])
    cf_paras = slice_to_next_marker(lines, cf_idx, markers) if cf_idx is not None else []

    # Synthesis
    syn_idx = find_heading_index(lines, [re.compile(r"2\.2\s+Synthesis", re.I), re.compile(r"Synthesis", re.I)])
    syn_paras = slice_to_next_marker(lines, syn_idx, markers) if syn_idx is not None else []

    # Gap (heuristic: gather paragraphs mentioning gap/weakness/limitations)
    gap_paras = [t for t in lines if re.search(r"\bgap\b|weakness|limitations", t, re.I)]

    # Methodology (Chapter 3 heading)
    meth_idx = find_heading_index(lines, [re.compile(r"METHODOLOGY", re.I)])
    meth_paras = slice_to_next_marker(lines, meth_idx, markers) if meth_idx is not None else []

    # Requirement analysis
    req_idx = find_heading_index(lines, [re.compile(r"3\.1\s+Requirement Analysis", re.I)])
    req_paras = slice_to_next_marker(lines, req_idx, markers) if req_idx is not None else []

    # Population and locality
    pop_idx = find_heading_index(lines, [re.compile(r"3\.1\.1\s+Population and Locale of the Study", re.I)])
    pop_paras = slice_to_next_marker(lines, pop_idx, markers) if pop_idx is not None else []

    # Data instrumentation
    inst_idx = find_heading_index(lines, [re.compile(r"3\.1\.2\s+Data Instrumentation", re.I)])
    inst_paras = slice_to_next_marker(lines, inst_idx, markers) if inst_idx is not None else []

    # Likert scale (heuristic by keyword)
    likert_paras = [t for t in lines if re.search(r"Likert", t, re.I)]

    # DFD / current technical solution: use Project Context paragraphs as "current".
    ctx_idx = find_heading_index(lines, [re.compile(r"Project Context", re.I)])
    ctx_paras = slice_to_next_marker(lines, ctx_idx, markers) if ctx_idx is not None else intro_paras

    # DFD proposed solution: section 3.5
    dfd_idx = find_heading_index(lines, [re.compile(r"3\.5\s+Tools for Data Analysis", re.I), re.compile(r"Data Flow Diagrams \(DFD\)", re.I)])
    dfd_paras = slice_to_next_marker(lines, dfd_idx, markers) if dfd_idx is not None else [t for t in lines if re.search(r"\bDFD\b|Data Flow", t, re.I)]

    # Prototype (heuristic from designing phase / framework)
    proto_paras = [t for t in lines if re.search(r"WPF|WebView2|prototype|UI|interface", t, re.I)]

    # System architecture / framework
    arch_idx = find_heading_index(lines, [re.compile(r"3\.4\s+System Framework", re.I), re.compile(r"System Framework", re.I)])
    arch_paras = slice_to_next_marker(lines, arch_idx, markers) if arch_idx is not None else [t for t in lines if re.search(r"architecture|framework", t, re.I)]

    # Implementation plan
    impl_idx = find_heading_index(lines, [re.compile(r"3\.6\s+The Proposed Implementation", re.I), re.compile(r"Proposed Implementation", re.I)])
    impl_paras = slice_to_next_marker(lines, impl_idx, markers) if impl_idx is not None else [t for t in lines if re.search(r"implementation", t, re.I)]

    # Statistical tools (will be merged into the Data Instrument slide)
    stats_idx = find_heading_index(lines, [re.compile(r"3\.7\s+Statistical Tools", re.I)])
    stats_paras = slice_to_next_marker(lines, stats_idx, markers) if stats_idx is not None else [t for t in lines if re.search(r"statistical|weighted arithmetic mean|mean", t, re.I)]

    title = extract_project_title(lines)

    # Build bullets
    intro_bullets_all = to_bullets(intro_paras, 14)
    intro_1, intro_2 = split_bullets(intro_bullets_all, 7)

    rol_bullets = to_bullets(rol_paras, 16)
    rol_1, rol_2 = split_bullets(rol_bullets, 8)

    prs = Presentation()

    add_title_slide(prs, title, members=[])

    add_bullets_slide(prs, "Introduction (1/2)", intro_1)
    add_bullets_slide(prs, "Introduction (2/2)", intro_2)

    add_bullets_slide(prs, "Objectives", to_bullets(obj_paras, 10))

    add_bullets_slide(prs, "Scope and Delimitation", to_bullets(scope_paras, 10))

    add_bullets_slide(prs, "Summary of Review of Related Literature", rol_1)
    add_bullets_slide(prs, "Conceptual Framework", to_bullets(cf_paras, 10))
    add_bullets_slide(prs, "Summary of Review of Related Studies", rol_2)

    add_bullets_slide(prs, "Synthesis", to_bullets(syn_paras, 10))
    add_bullets_slide(prs, "Research Gap", to_bullets(gap_paras, 8))

    add_bullets_slide(prs, "Methodology", to_bullets(meth_paras, 10))
    add_bullets_slide(prs, "Requirement Analysis", to_bullets(req_paras, 10))

    add_bullets_slide(prs, "DFD: Current Technical Solution (As-Is)", to_bullets(ctx_paras, 10))
    add_bullets_slide(prs, "DFD: Proposed Solution (SmartHub)", to_bullets(dfd_paras, 10))

    add_bullets_slide(prs, "Population and Locality", to_bullets(pop_paras, 10))
    add_bullets_slide(
        prs,
        "Data Instrument (Statistical Tool)",
        to_bullets(list(inst_paras) + list(stats_paras), 10),
    )
    add_bullets_slide(prs, "Likert Scale", to_bullets(likert_paras, 8))

    add_bullets_slide(prs, "Prototype", to_bullets(proto_paras, 10))
    add_bullets_slide(prs, "System Architecture", to_bullets(arch_paras, 10))
    add_bullets_slide(prs, "Implementation Plan", to_bullets(impl_paras, 10))

    OUT_PPTX_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PPTX_PATH))

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
